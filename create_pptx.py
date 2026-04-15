from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

PYPATH = r"C:\Users\bio\Desktop\python\RANSAC_ICP_v2.pptx"

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

FONT = "Meiryo"

# Colors
NAVY     = RGBColor(0x1A, 0x2E, 0x4A)
BLUE     = RGBColor(0x1D, 0x5F, 0xA5)
LBLUE    = RGBColor(0xD6, 0xE8, 0xFA)
TEAL     = RGBColor(0x0D, 0x94, 0x88)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF5, 0xF8, 0xFC)
DARK     = RGBColor(0x1E, 0x29, 0x3B)
MID      = RGBColor(0x47, 0x55, 0x69)
LIGHT    = RGBColor(0x94, 0xA3, 0xB8)
GOLD     = RGBColor(0xF5, 0x9E, 0x0B)
GREEN    = RGBColor(0x16, 0xA3, 0x4A)
RED      = RGBColor(0xDC, 0x26, 0x26)
DARKBLUE = RGBColor(0x1E, 0x3A, 0x5F)

blank_layout = prs.slide_layouts[6]  # blank

def i(val): return Inches(val)

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        i(x), i(y), i(w), i(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        9,  # MSO_SHAPE_TYPE.OVAL
        i(x), i(y), i(w), i(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(i(x), i(y), i(w), i(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

TOTAL = 9

def add_slide_num(slide, num):
    add_text(slide, f"{num} / {TOTAL}", 9.3, 5.25, 0.6, 0.25,
             size=10, color=LIGHT, align=PP_ALIGN.RIGHT)

def add_header(slide, title, color=None):
    c = color or NAVY
    add_rect(slide, 0, 0, 10, 1.0, c)
    add_text(slide, title, 0.5, 0.12, 9, 0.76,
             size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# =============================================
# SLIDE 1: Title
# =============================================
s1 = prs.slides.add_slide(blank_layout)
# Background
add_rect(s1, 0, 0, 10, 5.625, NAVY)
# Accent bars
add_rect(s1, 0, 0, 10, 0.12, TEAL)
add_rect(s1, 0, 5.5, 10, 0.125, TEAL)
# Decorative circle
add_oval(s1, 7.3, 0.7, 3.5, 3.5, RGBColor(0x1D, 0x4E, 0x89))
# Title
add_text(s1, "3Dモデル位置合わせの仕組み", 0.6, 1.4, 7.5, 1.3,
         size=34, bold=True, color=WHITE)
add_text(s1, "RANSAC・ICP とスケール補正", 0.6, 2.9, 7.0, 0.7,
         size=20, color=LBLUE)
add_slide_num(s1, 1)

# =============================================
# SLIDE 2: 処理の全体フロー
# =============================================
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, 10, 5.625, OFFWHITE)
add_header(s2, "処理の全体フロー")

steps = [
    ("1", "主軸アライメント\n（PCA）", "重心と主軸方向を\n一致させる前処理", TEAL),
    ("2", "RANSAC", "特徴点マッチングによる\n粗い位置合わせ", BLUE),
    ("3", "ICP", "点群全体を使った\n精密な位置合わせ\n・スケール補正", NAVY),
]

bw, bh = 2.4, 2.8
sx, sy, gap = 0.8, 1.2, 0.8

for idx, (num, label, desc, color) in enumerate(steps):
    x = sx + idx * (bw + gap)
    # Card
    add_rect(s2, x, sy, bw, bh, WHITE, color, Pt(2))
    # Colored top
    add_rect(s2, x, sy, bw, 0.55, color)
    # Step number
    add_oval(s2, x + 0.06, sy + 0.06, 0.43, 0.43, WHITE)
    add_text(s2, num, x + 0.06, sy + 0.07, 0.43, 0.41,
             size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Label
    add_text(s2, label, x + 0.1, sy + 0.58, bw - 0.2, 0.9,
             size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    # Desc
    add_text(s2, desc, x + 0.1, sy + 1.55, bw - 0.2, 1.2,
             size=12, color=MID, align=PP_ALIGN.CENTER)
    # Arrow
    if idx < 2:
        ax = x + bw + 0.25
        add_text(s2, "▶", ax, sy + bh / 2 - 0.2, 0.3, 0.4,
                 size=18, color=MID, align=PP_ALIGN.CENTER)

add_slide_num(s2, 2)

# =============================================
# SLIDE 3: RANSACの仕組み
# =============================================
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, 10, 5.625, OFFWHITE)
add_header(s3, "RANSAC（粗い位置合わせ）", BLUE)

add_rect(s3, 0.5, 1.1, 5.2, 0.4, LBLUE, BLUE, Pt(1))
add_text(s3, "Random Sample Consensus", 0.5, 1.12, 5.2, 0.36,
         size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

ransac_steps = [
    ("1", "特徴量計算", "点群から FPFH（Fast Point Feature Histograms）を計算"),
    ("2", "ランダムサンプリング", "特徴量が類似する点のペアをランダムに選択"),
    ("3", "変換行列の推定", "選択したペアから回転・並進・スケールを計算"),
    ("4", "Inlier 評価", "全点との対応数（Inlier数）でモデルを評価"),
    ("5", "最良解を採用", "最も多くの Inlier を持つ変換を最終解とする"),
]

for j, (n, t, d) in enumerate(ransac_steps):
    y = 1.65 + j * 0.68
    add_oval(s3, 0.5, y, 0.42, 0.42, BLUE)
    add_text(s3, n, 0.5, y + 0.02, 0.42, 0.38,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s3, t, 1.08, y, 2.0, 0.42,
             size=12, bold=True, color=DARK)
    add_text(s3, d, 3.2, y, 4.3, 0.42,
             size=11, color=MID)

# Feature box
add_rect(s3, 7.7, 1.1, 2.1, 3.55, LBLUE, BLUE, Pt(1))
add_text(s3, "特徴", 7.7, 1.1, 2.1, 0.38,
         size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
feats = ["✓ 外れ値（誤対応）に強い",
         "✓ 初期位置によらない\n  大域的探索が可能",
         "✓ 大まかな姿勢の\n  初期値を提供"]
for j, f in enumerate(feats):
    add_text(s3, f, 7.8, 1.58 + j * 1.0, 1.9, 0.9, size=11, color=DARK)

add_slide_num(s3, 3)

# =============================================
# SLIDE 4: ICPの仕組み
# =============================================
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, 10, 5.625, OFFWHITE)
add_header(s4, "ICP（精密な位置合わせ）", NAVY)

add_rect(s4, 0.5, 1.1, 4.5, 0.4, LBLUE, NAVY, Pt(1))
add_text(s4, "Iterative Closest Point", 0.5, 1.12, 4.5, 0.36,
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

icp_steps = [
    ("1", "最近傍点の対応付け", "各点に対して相手点群の最も近い点を対応付け"),
    ("2", "変換行列の計算", "対応点間の距離を最小化する変換（回転・並進・スケール）を計算"),
    ("3", "変換の適用", "計算した変換行列を点群に適用"),
    ("4", "収束まで繰り返す", "変化量が閾値を下回るか最大反復回数に達するまで1〜3を繰り返す"),
]

for j, (n, t, d) in enumerate(icp_steps):
    y = 1.65 + j * 0.82
    add_oval(s4, 0.5, y, 0.42, 0.42, NAVY)
    add_text(s4, n, 0.5, y + 0.02, 0.42, 0.38,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, t, 1.08, y, 2.2, 0.42,
             size=12, bold=True, color=DARK)
    add_text(s4, d, 3.4, y, 4.1, 0.42,
             size=11, color=MID)

add_rect(s4, 7.7, 1.1, 2.1, 3.55, LBLUE, NAVY, Pt(1))
add_text(s4, "特徴", 7.7, 1.1, 2.1, 0.38,
         size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
icp_feats = ["✓ 高精度なスケール補正",
             "✓ 点群全体を使った\n  局所最適化",
             "✓ RANSACの初期値で\n  高精度化"]
for j, f in enumerate(icp_feats):
    add_text(s4, f, 7.8, 1.58 + j * 1.0, 1.9, 0.9, size=11, color=DARK)

add_slide_num(s4, 4)

# =============================================
# SLIDE 5: なぜRANSACではスケール補正が難しいか
# =============================================
s5 = prs.slides.add_slide(blank_layout)
add_rect(s5, 0, 0, 10, 5.625, OFFWHITE)
add_header(s5, "RANSACでスケール補正が難しい理由", BLUE)

reasons = [
    ("①", "特徴点の対応付けが前提",
     ["FPFH特徴量の類似度で点を対応付ける",
      "スケール差が大きいと同じ部位でも特徴量が変化",
      "正しい対応が取れず推定が不正確になる"],
     BLUE),
    ("②", "ダウンサンプリングによる点群の粗さ",
     ["計算速度のため点群を間引いて使用",
      "疎な点群ではスケール変化を検出しにくい",
      "細かい形状の違いが失われる"],
     TEAL),
    ("③", "変換の自由度と最適化の難しさ",
     ["スケールを加えると変換行列の自由度が増加",
      "ランダムサンプリングで正しい\nスケールを引き当てるのが困難",
      "外れ値が多いとスケール推定がぶれやすい"],
     NAVY),
]

cw, ch = 2.9, 3.2
cx0, cy, gap = 0.35, 1.15, 0.2

for k, (num, title, pts, color) in enumerate(reasons):
    x = cx0 + k * (cw + gap)
    add_rect(s5, x, cy, cw, ch, WHITE, color, Pt(2))
    add_rect(s5, x, cy, cw, 0.55, color)
    add_text(s5, f"理由 {num}  {title}", x + 0.1, cy + 0.06, cw - 0.2, 0.44,
             size=11, bold=True, color=WHITE)
    for m, pt in enumerate(pts):
        add_text(s5, f"・ {pt}", x + 0.15, cy + 0.65 + m * 0.83, cw - 0.25, 0.78,
                 size=11, color=DARK)

# Conclusion
add_rect(s5, 0.35, 4.55, 9.3, 0.65, NAVY)
add_text(s5, "結論：RANSACの役割は「大まかな姿勢の初期値を与えること」。精密なスケール補正はICPが担う。",
         0.5, 4.57, 9.0, 0.6,
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_slide_num(s5, 5)

# =============================================
# SLIDE 6: 実際の処理結果の解釈
# =============================================
s6 = prs.slides.add_slide(blank_layout)
add_rect(s6, 0, 0, 10, 5.625, OFFWHITE)
add_header(s6, "実際の処理結果の解釈", NAVY)

stages = [
    ("処理前",   RGBColor(0x25,0x63,0xEB), RGBColor(0xDC,0x26,0x26), "2つのモデルが\n離れた位置にある"),
    ("RANSAC後", RGBColor(0xDC,0x26,0x26), RGBColor(0xEA,0xB3,0x08), "形状の重なりは改善\nスケール差が残る"),
    ("ICP後",    RGBColor(0xDC,0x26,0x26), RGBColor(0x16,0xA3,0x4A), "スケールを含めて\n精密に一致"),
]

for idx, (label, c1, c2, desc) in enumerate(stages):
    x = 0.5 + idx * 3.2
    add_rect(s6, x, 1.1, 2.8, 0.38, NAVY)
    add_text(s6, label, x, 1.12, 2.8, 0.34,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Two overlapping ovals
    off = 0.38 if idx == 0 else 0.0
    vert = 0.25 if idx == 0 else 0.0
    add_oval(s6, x + 0.3 - off, 1.65, 1.1, 1.0, c1, c1)
    add_oval(s6, x + 0.8 + off, 1.65 + vert, 1.1, 1.0, c2, c2)
    add_text(s6, desc, x, 3.0, 2.8, 0.8,
             size=12, color=DARK, align=PP_ALIGN.CENTER)

# Arrows
for idx in range(1, 3):
    ax = 0.5 + idx * 3.2 - 0.35
    add_text(s6, "▶", ax, 1.95, 0.3, 0.35,
             size=18, color=MID, align=PP_ALIGN.CENTER)

# Summary
add_rect(s6, 0.5, 4.62, 9.0, 0.55, LBLUE, BLUE, Pt(1))
add_text(s6, "→ RANSACは「位置合わせの道筋をつける役割」として正常に機能している",
         0.65, 4.64, 8.7, 0.5,
         size=12, bold=True, color=NAVY)

add_slide_num(s6, 6)

# =============================================
# SLIDE 7: まとめ
# =============================================
s7 = prs.slides.add_slide(blank_layout)
add_rect(s7, 0, 0, 10, 5.625, NAVY)
add_rect(s7, 0, 0, 10, 0.12, TEAL)
add_rect(s7, 0, 5.5, 10, 0.125, TEAL)

add_text(s7, "まとめ", 0.5, 0.2, 9, 0.7,
         size=28, bold=True, color=WHITE)

summary = [
    ("RANSAC",         BLUE,  "特徴点ベースの大域的探索。外れ値に強く初期位置によらないが、スケール補正の精度は限定的。"),
    ("ICP",            TEAL,  "最近傍点ベースの局所最適化。RANSACの初期値があることで高精度なスケール補正を実現。"),
    ("2段階の\n組み合わせ", GOLD, "RANSACで粗く合わせ → ICPで精密に仕上げることで、スケール差のある3Dモデルの高精度な位置合わせを実現。"),
]

for k, (label, color, text) in enumerate(summary):
    y = 1.1 + k * 1.28
    add_rect(s7, 0.5, y, 2.2, 0.9, color)
    add_text(s7, label, 0.5, y + 0.05, 2.2, 0.8,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s7, 2.85, y, 6.8, 0.9, DARKBLUE, color, Pt(1))
    add_text(s7, text, 3.0, y + 0.05, 6.55, 0.8,
             size=12, color=LBLUE)

add_slide_num(s7, 7)

# =============================================
# SLIDE 8: Umeyamaアルゴリズム
# =============================================
s8 = prs.slides.add_slide(blank_layout)
add_rect(s8, 0, 0, 10, 5.625, OFFWHITE)
add_header(s8, "Umeyamaアルゴリズム（スケール補正の数学的基盤）", RGBColor(0x6D, 0x28, 0xD9))

PURPLE = RGBColor(0x6D, 0x28, 0xD9)
LPURP  = RGBColor(0xED, 0xE9, 0xFE)

# What is it
add_rect(s8, 0.4, 1.1, 9.2, 0.48, LPURP, PURPLE, Pt(1))
add_text(s8, "対応点が既知のとき、スケール・回転・並進を一括で求める閉形式（解析解）の手法　[Umeyama 1991]",
         0.55, 1.13, 8.9, 0.42, size=12, bold=True, color=PURPLE)

# Formula area
add_rect(s8, 0.4, 1.68, 4.3, 2.3, WHITE, PURPLE, Pt(1))
add_text(s8, "最小化する目的関数", 0.5, 1.72, 4.1, 0.35, size=12, bold=True, color=PURPLE)
add_text(s8, "1/n  \u03a3  || yi - (s\u00b7R\u00b7xi + t) ||^2",
         0.6, 2.12, 4.0, 0.5, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(s8, "s : スケール（拡大縮小）\nR : 回転行列\nt : 並進ベクトル\nxi : 変換元の点\nyi : 変換先の点",
         0.6, 2.68, 3.9, 1.2, size=12, color=MID)

# How it works
add_rect(s8, 4.85, 1.68, 4.75, 2.3, WHITE, PURPLE, Pt(1))
add_text(s8, "解法の手順（解析解）", 4.95, 1.72, 4.5, 0.35, size=12, bold=True, color=PURPLE)
steps_u = [
    "1. 両点群の重心を計算・除去",
    "2. 共分散行列を計算",
    "3. SVD（特異値分解）で回転Rを求める",
    "4. 分散比からスケールsを計算",
    "5. 並進tを計算",
]
for j, st in enumerate(steps_u):
    add_text(s8, st, 4.95, 2.1 + j * 0.37, 4.5, 0.35, size=11, color=DARK)

# Key point: relationship with RANSAC/ICP
add_rect(s8, 0.4, 4.08, 9.2, 1.15, NAVY)
add_text(s8, "RANSACとICPにおけるUmeyamaの役割", 0.55, 4.1, 9.0, 0.35,
         size=12, bold=True, color=TEAL)

roles = [
    ("RANSAC", BLUE, "3点のランダムペアに適用 → 対応が誤りやすくスケール推定が不安定"),
    ("ICP",    TEAL, "全最近傍点ペア（数千点）に反復適用 → 密な対応でスケールが正確に収束"),
]
for k, (lbl, col, txt) in enumerate(roles):
    x = 0.55 + k * 4.6
    add_rect(s8, x, 4.5, 1.1, 0.55, col)
    add_text(s8, lbl, x, 4.52, 1.1, 0.51, size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s8, txt, x + 1.2, 4.52, 3.25, 0.51, size=11, color=LBLUE)

add_slide_num(s8, 8)

# =============================================
# SLIDE 9: まとめ（旧SLIDE7を繰り上げ）
# =============================================
s9 = prs.slides.add_slide(blank_layout)
add_rect(s9, 0, 0, 10, 5.625, NAVY)
add_rect(s9, 0, 0, 10, 0.12, TEAL)
add_rect(s9, 0, 5.5, 10, 0.125, TEAL)

add_text(s9, "まとめ", 0.5, 0.2, 9, 0.7,
         size=28, bold=True, color=WHITE)

summary2 = [
    ("RANSAC",         BLUE,   "特徴点ベースの大域的探索。外れ値に強いが、Umeyamaによるスケール補正の精度は限定的。"),
    ("ICP",            TEAL,   "最近傍点ベースの局所最適化。Umeyamaを反復適用することで高精度なスケール補正を実現。"),
    ("Umeyama",        RGBColor(0x6D,0x28,0xD9), "スケール・回転・並進を一括で求める閉形式解法。RANSACとICP両方の内部で使用される。"),
    ("2段階の\n組み合わせ", GOLD, "RANSACで粗く合わせ → ICPで精密に仕上げ。スケール差のある3Dモデルの高精度位置合わせを実現。"),
]

for k, (label, color, text) in enumerate(summary2):
    y = 0.95 + k * 1.1
    add_rect(s9, 0.5, y, 2.0, 0.85, color)
    add_text(s9, label, 0.5, y + 0.05, 2.0, 0.75,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s9, 2.65, y, 7.0, 0.85, DARKBLUE, color, Pt(1))
    add_text(s9, text, 2.8, y + 0.08, 6.75, 0.7,
             size=11, color=LBLUE)

add_slide_num(s9, 9)

# Save
prs.save(PYPATH)
print(f"Saved: {PYPATH}")
